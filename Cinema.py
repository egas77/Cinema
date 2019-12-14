from docx import Document
from pptx import Presentation
import xlsxwriter

import sys
import datetime
import sqlite3
import numpy as np
import pickle
from PyQt5.QtWidgets import QApplication, QMainWindow, QDialog, QTableWidgetItem, QFileDialog
from PyQt5.QtCore import Qt
from PyQt5.QtGui import QPixmap
from PIL import Image, ImageQt, ImageDraw, ImageFont

from UI.ui_main import Ui_MainWindow
from UI.ui_add_session_dialog import Ui_AddSessionDialog
from UI.ui_add_cinema_dialog import Ui_AddCinemaDialog
from UI.ui_add_hall_dialog import Ui_AddHallDialog
from UI.ui_add_film_dialog import Ui_AdFilmDialog
from UI.ui_details_session_dialog import Ui_DetalisSessionDialog


class Cinema(QMainWindow, Ui_MainWindow):
    def __init__(self):
        super().__init__()
        self.setupUi(self)
        self.load_data_base()

        self.open_session_btn.clicked.connect(self.open_session)
        self.remove_session_btn.clicked.connect(self.remove_session)
        self.add_sesion_btn.clicked.connect(self.add_session)

        self.sessions_table.cellClicked.connect(self.cell_click)

        self.action_docx.triggered.connect(self.export_to_docx)
        self.action_pptx.triggered.connect(self.export_to_pptx)
        self.action_xlsx.triggered.connect(self.export_to_xlsx)

    def export_to_docx(self):
        file = QFileDialog.getSaveFileName(self, 'Сохранить', '', '''
                                            Docx (*.docx);;
                                            Doc (*.doc)''')[0]
        if file:
            con = sqlite3.connect(data_base_path)
            cur = con.cursor()

            data = cur.execute(
                '''SELECT film, date, time, hall_name, cinema FROM sessions
                ORDER BY date, time'''
            ).fetchall()
            if data:
                today_date = datetime.date.today()
                last_month_data = []
                for session in data:
                    film_id = session[0]
                    date_iso = session[1]
                    time_iso = session[2]
                    hall_name = session[3]
                    cinema_id = session[4]

                    date = datetime.date.fromisoformat(date_iso)

                    if today_date >= date:
                        date_delta = today_date - date
                        if date_delta.days < 30:
                            film_name = cur.execute(
                                '''SELECT name FROM films WHERE id = ?''', (film_id,)
                            ).fetchone()[0]

                            cinema_name = cur.execute(
                                '''SELECT name FROM cinemas WHERE id = ?''', (cinema_id,)
                            ).fetchone()[0]

                            current_session = (
                                film_name, cinema_name, hall_name, date_iso, time_iso
                            )
                            last_month_data.append(current_session)

                document = Document()

                if last_month_data:
                    rows = len(last_month_data)
                    cols = len(last_month_data[0])

                    table = document.add_table(rows + 1, cols)

                    head_table = table.rows[0].cells
                    head_table[0].text = 'Название фильма'
                    head_table[1].text = 'Название кинотеатра'
                    head_table[2].text = 'Название зала'
                    head_table[3].text = 'Дата'
                    head_table[4].text = 'Время'
                    for count_session in range(len(last_month_data)):
                        film_name = last_month_data[count_session][0]
                        cinema_name = last_month_data[count_session][1]
                        hall_name = last_month_data[count_session][2]
                        date = last_month_data[count_session][3]
                        time = last_month_data[count_session][4]

                        current_row = table.rows[count_session + 1].cells

                        current_row[0].text = film_name
                        current_row[1].text = cinema_name
                        current_row[2].text = hall_name
                        current_row[3].text = date
                        current_row[4].text = time

                document.save(file)

    def export_to_pptx(self):
        file = QFileDialog.getSaveFileName(self, 'Сохранить', '',
                                           '''Pptx (*.pptx)''')[0]
        if file:
            prs = Presentation()

            con = sqlite3.connect(data_base_path)
            cur = con.cursor()

            films = cur.execute(
                '''SELECT name, description FROM films'''
            ).fetchall()

            names_films = list(map(lambda data: data[0], films))
            descriptions = list(map(lambda data: data[1], films))

            slide_layout = prs.slide_layouts[1]

            for number_film, names_film, description in zip(range(len(films)), names_films,
                                                            descriptions):
                slide = prs.slides.add_slide(slide_layout)

                shapes = slide.shapes
                title_shape = shapes.placeholders[0]
                body_shape = shapes.placeholders[1]

                title_shape.text = names_film
                body_shape.text = description

            prs.save(file)

    def export_to_xlsx(self):
        file = QFileDialog.getSaveFileName(self, 'Сохранить', '',
                                           '''
                                           Xlsx (*.xlsx);;
                                           Xls (*xls)
                                           ''')[0]
        if file:
            workbook = xlsxwriter.Workbook(file)
            worksheet = workbook.add_worksheet()
            chart = workbook.add_chart({'type': 'column'})

            con = sqlite3.connect(data_base_path)
            cur = con.cursor()

            cinemas = cur.execute(
                '''SELECT name FROM cinemas'''
            ).fetchall()
            cinemas = list(map(lambda data: data[0], cinemas))

            for cinema_name, number_cinema in zip(cinemas, range(len(cinemas))):
                hours_dict = dict.fromkeys(
                    ['00', '01', '02', '03', '04', '05', '06', '07', '08', '09', '10', '11', '12',
                     '13',
                     '14', '15', '16', '17', '18', '19', '20', '21', '22', '23'], 0
                )

                sessions_current_cinema = cur.execute(
                    '''SELECT hall, time FROM sessions WHERE
                cinema = (SELECT id FROM cinemas WHERE name = ?)''', (cinema_name,)
                ).fetchall()

                for session in sessions_current_cinema:
                    hall_bytes = session[0]
                    time_iso = session[1]

                    hall_array = pickle.loads(hall_bytes)
                    buy_ticket_count = np.sum(hall_array)

                    hour = time_iso.split(':')[0]

                    hours_dict[hour] += buy_ticket_count

                values = list(map(lambda data: data, hours_dict.values()))
                worksheet.write(number_cinema, 0, cinema_name)
                worksheet.write_row(number_cinema, 1, values)
                chart.add_series({
                    'values': ['Sheet1', number_cinema, 1, number_cinema, 22],
                    'name': ['Sheet1', number_cinema, 0],
                })

            worksheet.insert_chart(len(cinemas), 0, chart)
            workbook.close()

    def load_data_base(self):
        con = sqlite3.connect(data_base_path)
        cur = con.cursor()

        sessions = cur.execute(
            '''SELECT film, date, time, hall_name, cinema FROM sessions'''
        ).fetchall()

        if sessions:
            self.sessions_table.clear()
            self.sessions_table.setRowCount(len(sessions))
            self.sessions_table.setColumnCount(5)

            head = ['Фильм', 'Кинотеатр', 'Зал', 'Дата', 'Время']
            self.sessions_table.setHorizontalHeaderLabels(head)
            for session_number in range(len(sessions)):
                film_id, date_iso, time_iso, hall_name, cinema = sessions[session_number]

                film_name = cur.execute(
                    '''SELECT name FROM films WHERE id = ?''', (film_id,)
                ).fetchall()[0][0]

                cinema_name = cur.execute(
                    '''SELECT name FROM cinemas WHERE id = ?''', (cinema,)
                ).fetchall()[0][0]

                self.sessions_table.setItem(session_number, 0, QTableWidgetItem(film_name))
                self.sessions_table.setItem(session_number, 1, QTableWidgetItem(cinema_name))
                self.sessions_table.setItem(session_number, 2, QTableWidgetItem(hall_name))
                self.sessions_table.setItem(session_number, 3, QTableWidgetItem(date_iso))
                self.sessions_table.setItem(session_number, 4, QTableWidgetItem(time_iso))

            self.sessions_table.resizeColumnsToContents()

        cur.close()
        con.close()

    def open_session(self):
        row = self.sessions_table.currentRow()

        if row != -1:
            film_name = self.sessions_table.item(row, 0).text()
            cinema_name = self.sessions_table.item(row, 1).text()
            hall_name = self.sessions_table.item(row, 2).text()
            date = self.sessions_table.item(row, 3).text()
            time = self.sessions_table.item(row, 4).text()

            con = sqlite3.connect(data_base_path)
            cur = con.cursor()

            id_session = cur.execute(
                '''SELECT id FROM sessions WHERE
                film = (SELECT id FROM films WHERE name = ?) AND
                date = ? AND time = ? AND
                cinema = (SELECT id FROM cinemas WHERE name = ?) AND
                hall_name = ?''',
                (film_name, date, time, cinema_name, hall_name,)).fetchone()[0]

            cur.close()
            con.close()

            self.details_session_dialog = DetailsSessionDialog(id_session)
            self.details_session_dialog.show()

    def remove_session(self):
        row = self.sessions_table.currentRow()

        if row != -1:
            film_name = self.sessions_table.item(row, 0).text()
            cinema_name = self.sessions_table.item(row, 1).text()
            hall_name = self.sessions_table.item(row, 2).text()
            date = self.sessions_table.item(row, 3).text()
            time = self.sessions_table.item(row, 4).text()

            con = sqlite3.connect(data_base_path)
            cur = con.cursor()

            id_session = cur.execute(
                '''SELECT id FROM sessions WHERE
                film = (SELECT id FROM films WHERE name = ?) AND
                date = ? AND time = ? AND
                cinema = (SELECT id FROM cinemas WHERE name = ?) AND
                hall_name = ?''',
                (film_name, date, time, cinema_name, hall_name,)).fetchone()[0]

            cur.execute('DELETE FROM sessions WHERE id = ?', (id_session,))

            con.commit()
            cur.close()
            con.close()

            self.load_data_base()

    def add_session(self):
        add_session_dialog = AddSessionDialog()
        add_session_dialog.accepted.connect(self.load_data_base)
        add_session_dialog.exec()

    def cell_click(self, row):
        self.sessions_table.selectRow(row)


class AddSessionDialog(QDialog, Ui_AddSessionDialog):
    def __init__(self):
        super().__init__()
        self.setupUi(self)
        self.setWindowFlag(Qt.WindowContextHelpButtonHint, False)
        self.setWindowTitle('Добавить сеанс')

        self.button_box.clicked.connect(self.click_button_box)
        self.cinema_tool_btn.clicked.connect(self.add_cinema)
        self.hall_tool_btn.clicked.connect(self.add_hall)
        self.film_tool_btn.clicked.connect(self.add_film)

        self.init_cinemas()
        self.init_halls()
        self.init_films()

        self.cinema_combo_box.currentTextChanged.connect(self.init_halls)

        self.data_edit.setDate(datetime.date.today())

    def init_cinemas(self):
        con = sqlite3.connect(data_base_path)
        cur = con.cursor()

        self.cinema_combo_box.clear()
        cinemas = cur.execute('SELECT name FROM cinemas').fetchall()
        if cinemas:
            cinemas = list(map(lambda data: data[0], cinemas))
            self.cinema_combo_box.addItems(cinemas)
            self.cinema_combo_box.setCurrentIndex(0)

        cur.close()
        con.close()

    def init_halls(self):
        self.hall_combo_box.clear()

        name_cinema = self.cinema_combo_box.currentText()
        if name_cinema:
            con = sqlite3.connect(data_base_path)
            cur = con.cursor()

            halls = cur.execute(
                'SELECT name FROM halls WHERE cinema = (SELECT id FROM cinemas WHERE name = ?)',
                (name_cinema,)
            ).fetchall()

            if halls:
                halls = list(map(lambda data: data[0], halls))
                self.hall_combo_box.addItems(halls)

            cur.close()
            con.close()

    def init_films(self):
        self.film_combo_box.clear()
        con = sqlite3.connect(data_base_path)
        cur = con.cursor()
        films = cur.execute('SELECT name FROM films').fetchall()
        if films:
            films = list(map(lambda data: data[0], films))
            self.film_combo_box.addItems(films)
        cur.close()
        con.close()

    def click_button_box(self, btn):
        btn_text = btn.text()
        if btn_text == 'OK':
            cinema_name = self.cinema_combo_box.currentText()
            hall_name = self.hall_combo_box.currentText()
            film_name = self.film_combo_box.currentText()
            date = self.data_edit.date().toPyDate().isoformat()
            time = self.time_edit.time().toPyTime().isoformat()

            if cinema_name and hall_name and film_name and date and time:
                con = sqlite3.connect(data_base_path)
                cur = con.cursor()

                cinema = cur.execute(
                    '''SELECT id FROM cinemas WHERE name = ?''', (cinema_name,)
                ).fetchall()[0][0]

                hall = cur.execute(
                    '''SELECT hall FROM halls WHERE
                     cinema = (SELECT id FROM cinemas WHERE name = ?) AND
                     name = ?''', (cinema_name, hall_name,)
                ).fetchall()[0][0]

                film_id = cur.execute(
                    '''SELECT id FROM films WHERE name = ?''', (film_name,)
                ).fetchall()[0][0]

                if not cur.execute('''SELECT * FROM sessions WHERE
                                    film = ? AND
                                    date = ? AND
                                    time = ? AND
                                    hall = ? AND
                                    hall_name = ? AND
                                    cinema = ?
                ''',
                                   (film_id, date, time, hall, hall_name, cinema,)).fetchall():
                    cur.execute(
                        '''INSERT INTO sessions(film, date, time, hall, hall_name, cinema)
                         VALUES (?, ?, ?, ?, ?, ?)''',
                        (film_id, date, time, hall, hall_name, cinema,)
                    )

                    con.commit()
                    cur.close()
                    con.close()

                    self.accept()
        elif btn_text == 'Cancel':
            self.reject()

    def add_cinema(self):
        add_cinema_dialog = AddCinemaDialog()
        add_cinema_dialog.accepted.connect(self.init_cinemas)
        add_cinema_dialog.exec()

    def add_hall(self):
        current_cinema = self.cinema_combo_box.currentText()
        if current_cinema:
            add_hall_dialog = AddHallDialog(current_cinema)
        else:
            add_hall_dialog = AddHallDialog()
        add_hall_dialog.accepted.connect(self.init_halls)
        add_hall_dialog.exec()

    def add_film(self):
        add_film_dialog = AddFilmDialog()
        add_film_dialog.accepted.connect(self.init_films)
        add_film_dialog.exec()


class AddCinemaDialog(QDialog, Ui_AddCinemaDialog):
    def __init__(self):
        super().__init__()
        self.setupUi(self)
        self.button_box.clicked.connect(self.click_button_box)
        self.setWindowFlag(Qt.WindowContextHelpButtonHint, False)
        self.setWindowTitle('Кинотеатр')

    def click_button_box(self, btn):
        btn_text = btn.text()
        if btn_text == 'OK':
            name = self.name_line_edit.text()
            address = self.address_line_edit.text()
            if name and address:
                con = sqlite3.connect(data_base_path)
                cur = con.cursor()
                cur.execute(
                    '''INSERT INTO cinemas(name, address) VALUES (?, ?)''', (name, address,)
                )
                con.commit()
                cur.close()
                con.close()
                self.accept()
        elif btn_text == 'Cancel':
            self.reject()


class AddHallDialog(QDialog, Ui_AddHallDialog):
    def __init__(self, cinema=None):
        super().__init__()
        self.setupUi(self)
        self.setWindowFlag(Qt.WindowContextHelpButtonHint, False)

        self.cinema = cinema
        self.init_cinemas()
        self.button_box.clicked.connect(self.click_button_box)

    def click_button_box(self, btn):
        btn_text = btn.text()
        if btn_text == 'OK':
            cinema_name = self.cinema_combo_box.currentText()
            hall_name = self.name_line_edit.text()
            count_row = self.count_row_edit.value()
            count_col = self.count_col_elit.value()
            if cinema_name and hall_name and count_row and count_row:
                con = sqlite3.connect(data_base_path)
                cur = con.cursor()
                hall = np.full((count_row, count_col), 0, dtype=int)
                hall_bytes = hall.dumps()
                cur.execute(
                    '''INSERT INTO halls(cinema, hall, name) VALUES (
                    (SELECT id FROM cinemas WHERE name = ?), ?, ?)''',
                    (cinema_name, hall_bytes, hall_name)
                )
                con.commit()
                cur.close()
                con.close()
                self.accept()
        elif btn_text == 'Cancel':
            self.reject()

    def init_cinemas(self):
        con = sqlite3.connect(data_base_path)
        cur = con.cursor()

        self.cinema_combo_box.clear()
        cinemas = cur.execute('SELECT name FROM cinemas').fetchall()
        if cinemas:
            cinemas = list(map(lambda data: data[0], cinemas))
            self.cinema_combo_box.addItems(cinemas)
            if self.cinema:
                self.cinema_combo_box.setCurrentText(self.cinema)
            else:
                self.cinema_combo_box.setCurrentIndex(0)

        con.close()


class AddFilmDialog(QDialog, Ui_AdFilmDialog):
    def __init__(self):
        super().__init__()
        self.setupUi(self)
        self.setWindowTitle('Фильм')
        self.setWindowFlag(Qt.WindowContextHelpButtonHint, False)
        self.button_box.clicked.connect(self.click_button_box)

    def click_button_box(self, btn):
        btn_text = btn.text()
        if btn_text == 'OK':
            name = self.name_line_edit.text()
            description = self.description_edit.toPlainText()
            if name and description:
                con = sqlite3.connect(data_base_path)
                cur = con.cursor()
                cur.execute(
                    '''INSERT INTO films(name, description) VALUES (?, ?)''',
                    (name, description,)
                )
                con.commit()
                cur.close()
                con.close()
                self.accept()
        elif btn_text == 'Cancel':
            self.reject()


class DetailsSessionDialog(QDialog, Ui_DetalisSessionDialog):
    def __init__(self, id_session):
        super().__init__()
        self.setupUi(self)
        self.id_session = id_session
        self.size_font = 75
        self.setWindowFlag(Qt.WindowContextHelpButtonHint, False)
        self.buy_ticket_btn.clicked.connect(self.buy_ticket)

        self.initUI()

    def buy_ticket(self):
        row = self.row_edit.value()
        col = self.col_edit.value()

        con = sqlite3.connect(data_base_path)
        cur = con.cursor()

        hall_bytes = cur.execute(
            '''SELECT hall FROM sessions WHERE id = ?''', (self.id_session,)
        ).fetchone()[0]

        hall_array = pickle.loads(hall_bytes)
        shape = hall_array.shape

        if row <= shape[0] and col <= shape[1]:
            hall_array[row - 1][col - 1] = 1
            hall_bytes = hall_array.dumps()
            cur.execute(
                '''UPDATE sessions SET hall = ? WHERE id = ?''', (hall_bytes, self.id_session,)
            )
            con.commit()
            cur.close()
            con.close()
            self.resizeEvent(None)

    def initUI(self):
        con = sqlite3.connect(data_base_path)
        cur = con.cursor()

        session_data = cur.execute('''SELECT id, film, date, time, hall FROM sessions WHERE
        id = ?''', (self.id_session,)).fetchall()[0]

        name_film = cur.execute(
            '''SELECT name FROM films WHERE id = ?''', (session_data[1],)
        ).fetchone()[0]
        date = session_data[2]
        time = session_data[3]
        self.id_hall = session_data[4]

        self.name_film_label.setText(name_film)
        self.date_label.setText(date)
        self.time_label.setText(time)

    def resizeEvent(self, event):
        self.load_image_hall()

    def load_image_hall(self):
        con = sqlite3.connect(data_base_path)
        cur = con.cursor()

        hall_bytes = cur.execute(
            '''SELECT hall FROM sessions WHERE id = ?''', (self.id_session,)
        ).fetchone()[0]

        hall_array = pickle.loads(hall_bytes)
        font = ImageFont.truetype("arial.ttf", self.size_font)
        height_text = self.size_font
        max_width_text_row, max_width_text_col = 0, 0

        image = Image.new('RGB', (5000, 5000))
        draw = ImageDraw.Draw(image)

        shape = hall_array.shape

        for num in range(shape[0]):
            width_text, height_text = draw.textsize(str(num + 1), font=font)
            if width_text > max_width_text_row:
                max_width_text_row = width_text

        for num in range(shape[1]):
            width_text, height_text = draw.textsize(str(num + 1), font=font)
            if width_text > max_width_text_col:
                max_width_text_col = width_text

        for row in range(1, shape[0] + 1):
            draw.text((0, row * height_text), str(row), font=font)

        for col in range(shape[1]):
            draw.text((col * max_width_text_col + max_width_text_row, 0),
                      str(col + 1), font=font)

        size_img = (max_width_text_row + max_width_text_col * shape[1],
                    height_text + height_text * shape[0])

        for x, col in zip(range(max_width_text_row, size_img[0], max_width_text_col),
                          range(shape[1])):
            for y, row in zip(range(height_text, size_img[1], height_text), range(shape[0])):

                elem = hall_array[row][col]
                if elem == 1:
                    draw.rectangle([(x, y), (x + max_width_text_col - 1, y + height_text - 1)],
                                   width=2, fill=(255, 0, 0), outline=(255, 255, 255))
                else:
                    draw.rectangle([(x, y), (x + max_width_text_col - 1, y + height_text - 1)],
                                   width=2)

        image = image.crop((0, 0, size_img[0] + 1, size_img[1] + 1))

        image_qt = ImageQt.ImageQt(image)
        pixmap = QPixmap.fromImage(image_qt)
        pixmap = pixmap.scaled(self.img_cont.size(), Qt.KeepAspectRatio)
        self.img_cont.setPixmap(pixmap)


if __name__ == '__main__':
    data_base_path = 'data/cinema.db'
    app = QApplication(sys.argv)
    cinema = Cinema()
    cinema.show()
    sys.exit(app.exec())
